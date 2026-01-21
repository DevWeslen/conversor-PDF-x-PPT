import streamlit as st
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import tempfile
import os
import platform
import base64

# ===== CONFIG STREAMLIT =====
st.set_page_config(
    page_title="PDF → PPTX",
    layout="wide"
)

# ===== CARREGAR LOGO COMO BASE64 =====
logo_path = "PrincesadosCampos-positivo2-preferencial (1).jpg"
with open(logo_path, "rb") as f:
    logo_bytes = f.read()
logo_base64 = base64.b64encode(logo_bytes).decode()

# ===== ESTILO CSS =====
# ===== ESTILO CSS =====
st.markdown(
    f"""
    <style>
        /* Background da página */
        body, .main {{
            background-color: #f2f2f2;
        }}
        /* Cabeçalho com cor #043c1c, logo à esquerda e título centralizado */
        .custom-header {{
            background-color: #043c1c;
            color: white;
            display: flex;
            align-items: center;
            padding: 10px 20px;
            font-size: 28px;
            font-weight: bold;
            border-radius: 5px;
            margin-bottom: 20px;
        }}
        .custom-header img {{
            height: 100px; /* Aumentei aqui */
            margin-right: 15px;
        }}
        .custom-header span {{
            flex: 1;
            text-align: center;
        }}
    </style>
    """,
    unsafe_allow_html=True
)

# ===== CABEÇALHO =====
st.markdown(
    f"""
    <div class="custom-header">
        <img src="data:image/jpeg;base64,{logo_base64}" alt="Logo">
        <span>Conversor de PDF para PowerPoint </span>
    </div>
    """,
    unsafe_allow_html=True
)

st.write("Converte PDF mantendo layout fiel (imagem por slide)")

# ===== UPLOAD =====
uploaded_pdf = st.file_uploader(
    "Selecione o PDF",
    type=["pdf"]
)

if uploaded_pdf:
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, "entrada.pdf")
        with open(pdf_path, "wb") as f:
            f.write(uploaded_pdf.read())

        try:
            # Define poppler_path apenas no Windows
            poppler_dir = None
            if platform.system() == "Windows":
                poppler_dir = r"C:\poppler-25.12.0\Library\bin"

            pages = convert_from_path(
                pdf_path,
                dpi=150,
                poppler_path=poppler_dir
            )

            st.write("### 📄 Visualização das páginas do PDF:")
            page_selection = st.slider(
                "Escolha a página para visualizar",
                min_value=1,
                max_value=len(pages),
                value=1
            )

            # Prepara imagem para exibição com borda e centralizada
            img_bytes = BytesIO()
            pages[page_selection - 1].save(img_bytes, format="PNG")
            img_bytes.seek(0)
            encoded = base64.b64encode(img_bytes.read()).decode()

            st.markdown(
                f"""
                <div style="display: flex; justify-content: center; margin: 40px 0;">
                    <img src="data:image/png;base64,{encoded}" 
                         style="
                             border: 3px solid #000000; 
                             border-radius: 10px; 
                             max-width: 100%; 
                             height: auto;
                         ">
                </div>
                """,
                unsafe_allow_html=True
            )

            # Botão de conversão para PPTX
            if st.button("Converter todas as páginas para PPTX"):
                with st.spinner("Convertendo todas as páginas..."):
                    prs = Presentation()
                    prs.slide_width = Inches(13.33)
                    prs.slide_height = Inches(7.5)

                    for img in pages:
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        img_stream = BytesIO()
                        img.save(img_stream, format="PNG")
                        img_stream.seek(0)

                        slide.shapes.add_picture(
                            img_stream,
                            Inches(0),
                            Inches(0),
                            prs.slide_width,
                            prs.slide_height
                        )

                    ppt_path = os.path.join(tmpdir, "saida.pptx")
                    prs.save(ppt_path)

                    with open(ppt_path, "rb") as f:
                        st.success("✅ Conversão concluída!")
                        st.download_button(
                            "⬇️ Baixar PPTX",
                            f,
                            file_name="convertido.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

        except Exception as e:
            st.error(f"❌ Erro ao processar PDF: {e}")
