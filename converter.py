from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

PDF_PATH = "entrada.pdf"
PPT_PATH = "saida_final_imagem.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

paginas = convert_from_path(
    PDF_PATH,
    dpi=300,
    poppler_path=r"C:\poppler\Library\bin"
)

for i, img in enumerate(paginas):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    temp = f"page_{i}.png"
    img.save(temp)

    slide.shapes.add_picture(
        temp,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height
    )

    os.remove(temp)

prs.save(PPT_PATH)
print("✅ PPT gerado SOMENTE COM IMAGENS, fiel ao PDF.")
